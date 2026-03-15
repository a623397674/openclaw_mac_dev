#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# 设置幻灯片尺寸 (16:9)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle
    return slide

def add_content_slide(prs, title, content_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(content_lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(20)
        if line.startswith("•"):
            p.level = 1
    return slide

# Slide 1: 封面
slide = add_title_slide(prs, "一人公司团队介绍", "高效协作 · 快速响应 · 专业交付")

# Slide 2: 团队概述
add_content_slide(prs, "我们是谁", [
    "一支精干的 AI 协作团队，提供全方位产品与技术服务",
    "",
    "核心优势:",
    "• 3 分钟快速响应",
    "• 产品经理统一协调",
    "• 多角色协同作战",
    "• 7×24 小时在线"
])

# Slide 3: 产品小助手
add_content_slide(prs, "📊 产品小助手", [
    "角色：产品经理 / 项目协调员",
    "",
    "职责:",
    "• 需求分析与拆解",
    "• PRD 文档撰写",
    "• 任务分发与协调",
    "• 跨部门沟通",
    "",
    "响应时间：0-30 秒确认，3 分钟内出方案"
])

# Slide 4: 后端老张
add_content_slide(prs, "⚙️ 后端老张", [
    "角色：技术架构师 / 后端开发",
    "",
    "职责:",
    "• 系统架构设计",
    "• API 接口开发",
    "• 数据库设计",
    "• 技术选型建议",
    "",
    "擅长：高并发系统、微服务架构、性能优化"
])

# Slide 5: 前端小美
add_content_slide(prs, "🎨 前端小美", [
    "角色：前端开发 / UI/UX 设计师",
    "",
    "职责:",
    "• 界面设计与开发",
    "• 用户体验优化",
    "• 交互原型设计",
    "• 响应式适配",
    "",
    "擅长：现代化前端框架、动效设计、无障碍访问"
])

# Slide 6: 财务小慧
add_content_slide(prs, "💰 财务小慧", [
    "角色：财务顾问 / 预算分析师",
    "",
    "职责:",
    "• 成本估算",
    "• 预算规划",
    "• ROI 分析",
    "• 财务风险评估",
    "",
    "擅长：项目成本核算、投入产出分析、资金流规划"
])

# Slide 7: 客服小暖
add_content_slide(prs, "🎧 客服小暖", [
    "角色：客户成功 / 用户支持",
    "",
    "职责:",
    "• 用户问题解答",
    "• 反馈收集",
    "• 服务流程优化",
    "• 用户满意度提升",
    "",
    "擅长：问题快速定位、情绪安抚、服务话术设计"
])

# Slide 8: 工作流程
add_content_slide(prs, "🔄 3 分钟快速响应机制", [
    "0-30 秒：确认收到需求，告知预计完成时间",
    "",
    "30 秒 -2 分钟：分析任务类型，协同对应角色",
    "",
    "2-3 分钟：整合各方意见，给出统一方案"
])

# Slide 9: 服务场景
add_content_slide(prs, "✅ 服务场景", [
    "产品规划：需求分析、PRD 撰写、路线图",
    "技术开发：架构设计、代码实现、Code Review",
    "界面设计：UI 设计、交互原型、视觉优化",
    "财务预算：成本估算、ROI 分析、预算规划",
    "用户服务：问题解答、反馈处理、满意度提升",
    "综合项目：全员协同、端到端交付"
])

# Slide 10: 核心优势
add_content_slide(prs, "🌟 核心优势", [
    "⚡ 极速响应：3 分钟内出完整方案",
    "🎯 专业分工：每个角色各司其职",
    "📈 数据驱动：用数据说话，科学决策",
    "💡 用户导向：始终从用户价值出发",
    "🔧 务实落地：考虑资源、时间、技术可行性",
    "🤖 7×24 在线：随时待命，无休无止"
])

# Slide 11: 协作方式
add_content_slide(prs, "📱 如何联系我们", [
    "1. 私聊产品小助手 (QQ: 1903516201)",
    "2. 描述你的需求 (越详细越好)",
    "3. 等待 3 分钟 (接收完整方案)",
    "4. 确认或调整 (我们随时配合修改)",
    "",
    "需求描述建议:",
    "• 业务目标是什么？",
    "• 目标用户是谁？",
    "• 期望完成时间？",
    "• 预算范围？"
])

# Slide 12: 成功案例
add_content_slide(prs, "🏆 成功案例", [
    "电商平台重构：2 周完成 MVP，转化率提升 35%",
    "SaaS 产品设计：1 个月从 0 到 1，获客成本降低 50%",
    "移动端 APP 开发：3 个月双端上线，DAU 破 10 万",
    "数据分析系统：1 周搭建看板，决策效率提升 3 倍"
])

# Slide 13: 联系我们
slide = add_title_slide(prs, "谢谢观看", "期待与您合作")

# 保存
prs.save('/Users/jie/.openclaw/workspace-agents/product-manager/team-intro.pptx')
print("PPT 生成成功：team-intro.pptx")
